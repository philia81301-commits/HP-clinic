#!/usr/bin/env python3
"""Assemble the HP-clinic image deck PPTX from spec.yaml + slides/images/*.png."""
import argparse
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT_NAME = "jf-openhuninn-2.1"


def resolve_color(palette, token):
    hexval = palette.get(token, token)
    hexval = hexval.lstrip("#")
    return RGBColor.from_string(hexval.upper())


def add_overlay(slide, block, palette):
    left, top = Inches(block["x"]), Inches(block["y"])
    width, height = Inches(block["w"]), Inches(block["h"])
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    align = align_map.get(block.get("align", "left"), PP_ALIGN.LEFT)
    color = resolve_color(palette, block.get("color", "text"))
    weight = block.get("weight", "regular")
    font_pt = block.get("font_pt", 18)

    lines = str(block["text"]).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(font_pt)
        run.font.bold = weight == "bold"
        run.font.color.rgb = color


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True)
    parser.add_argument("--images-dir", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()

    spec = yaml.safe_load(Path(args.spec).read_text(encoding="utf-8"))
    palette = spec["design_system"]["palette"]
    size_in = spec["canvas"]["pptx_size_in"]
    images_dir = Path(args.images_dir)

    prs = Presentation()
    prs.slide_width = Emu(int(size_in["width"] * 914400))
    prs.slide_height = Emu(int(size_in["height"] * 914400))
    blank_layout = prs.slide_layouts[6]

    for slide_spec in spec["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        img_name = Path(slide_spec["output"]).name
        img_path = images_dir / img_name
        slide.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                  width=prs.slide_width, height=prs.slide_height)
        for block in slide_spec.get("overlay_blocks", []):
            add_overlay(slide, block, palette)

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved: {out_path} ({len(spec['slides'])} slides)")


if __name__ == "__main__":
    main()
